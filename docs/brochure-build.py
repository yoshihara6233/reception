"""
展示会用 製品紹介チラシ（1 枚裏表）を PPTX で生成する。

- A4 縦サイズ（210mm × 297mm）2 ページ
- 表面: ヒーロー + 主要4機能 + 価値訴求
- 裏面: 機能マトリクス + アーキ概略 + 料金 + 問合せ

実行:
  python3 docs/brochure-build.py
出力:
  docs/recorder-monitoring-brochure.pptx
"""
from pptx import Presentation
from pptx.util import Mm, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from copy import deepcopy
from lxml import etree

OUT = "docs/recorder-monitoring-brochure.pptx"

# ─── デザインシステム ────────────────────────────────────────────────
NAVY     = RGBColor(0x1E, 0x3A, 0x8A)   # primary
BLUE     = RGBColor(0x25, 0x63, 0xEB)   # accent
LIGHT    = RGBColor(0xEF, 0xF6, 0xFF)   # blue-50
INK      = RGBColor(0x0F, 0x17, 0x2A)   # slate-900
INK2     = RGBColor(0x33, 0x41, 0x55)   # slate-700
MUTED    = RGBColor(0x64, 0x74, 0x8B)   # slate-500
LINE     = RGBColor(0xCB, 0xD5, 0xE1)   # slate-300
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
GREEN    = RGBColor(0x15, 0x80, 0x3D)
AMBER    = RGBColor(0xB4, 0x53, 0x09)
RED      = RGBColor(0xB9, 0x1C, 0x1C)

FONT     = "Hiragino Kaku Gothic ProN"   # macOS 標準
FONT_BOLD = "Hiragino Kaku Gothic StdN W8"

# ─── プレゼンテーション初期化（A4 縦） ────────────────────────────────
prs = Presentation()
prs.slide_width  = Mm(210)
prs.slide_height = Mm(297)


def add_blank_slide():
    blank = prs.slide_layouts[6]  # blank layout
    return prs.slides.add_slide(blank)


def add_rect(slide, x, y, w, h, fill=None, line=None, line_w=0):
    """単色矩形。fill=None で透明。"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(line_w) if line_w else Pt(0.75)
    shape.shadow.inherit = False
    return shape


def add_text(
    slide, x, y, w, h, text,
    *, size=11, color=INK, bold=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
    font=FONT, line_spacing=1.25,
):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Mm(0)
    tf.margin_top = tf.margin_bottom = Mm(0)
    tf.vertical_anchor = anchor
    lines = text.split("\n")
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        r = p.add_run()
        r.text = ln
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
        # 日本語フォント明示
        rPr = r._r.get_or_add_rPr()
        ea = rPr.find(qn("a:ea"))
        if ea is None:
            ea = etree.SubElement(rPr, qn("a:ea"))
        ea.set("typeface", font)
    return tb


# ─────────────────────────────────────────────────────────────────────
# ▼ Slide 1 (表面)
# ─────────────────────────────────────────────────────────────────────
s1 = add_blank_slide()

# ヘッダーバー（ナビ色のフルブリード）
add_rect(s1, Mm(0), Mm(0), Mm(210), Mm(60), fill=NAVY)
add_rect(s1, Mm(0), Mm(60), Mm(210), Mm(3), fill=BLUE)

# 製品ロゴ風（左上）
add_text(
    s1, Mm(12), Mm(10), Mm(80), Mm(8),
    "RECORDER MONITOR  /  本部一元監視 SaaS",
    size=11, color=RGBColor(0xCB, 0xD5, 0xE1), font=FONT,
)
# メインタイトル
add_text(
    s1, Mm(12), Mm(18), Mm(186), Mm(20),
    "10,000 店舗を 1 名で運用できる、\nレコーダ統合監視プラットフォーム",
    size=22, color=WHITE, bold=True, font=FONT_BOLD, line_spacing=1.2,
)
# サブタイトル
add_text(
    s1, Mm(12), Mm(46), Mm(186), Mm(10),
    "i-PRO / Uniview / Frigate 対応  ×  JPEG ポーリング方式  ×  PWA 対応",
    size=11, color=RGBColor(0xDB, 0xEA, 0xFE), font=FONT,
)

# キーバッジ（白抜き）
def header_badge(x, label):
    add_rect(s1, Mm(x), Mm(64), Mm(40), Mm(7),
             fill=RGBColor(0xDB, 0xEA, 0xFE), line=BLUE, line_w=0.5)
    add_text(s1, Mm(x), Mm(64), Mm(40), Mm(7), label,
             size=8.5, color=NAVY, bold=True, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE, font=FONT_BOLD)

for i, lab in enumerate([
    "本部 1 名運用",
    "オンデマンド方式",
    "ポート開放不要",
    "PWA + Web Push",
]):
    header_badge(12 + i * 47, lab)

# キャッチコピー
add_text(
    s1, Mm(12), Mm(78), Mm(186), Mm(10),
    "「全店舗 24h ライブ配信」をやめれば、\nクラウドコストは 95% 削減できる。",
    size=15, color=INK, bold=True, font=FONT_BOLD, line_spacing=1.25,
)

# 説明文
add_text(
    s1, Mm(12), Mm(100), Mm(186), Mm(20),
    "視聴開始時にのみエッジが動く「オンデマンド方式」と、\n"
    "RTSP → ffmpeg → WebRTC の重いスタックを必要時のみに絞った「JPEG ポーリング方式」の組み合わせで、\n"
    "大規模化でも 1 店舗あたり月額 ¥58 のクラウドコストを実現。",
    size=10, color=INK2, font=FONT, line_spacing=1.5,
)

# ─── 主要 4 機能カード（2 x 2） ───────────────────────────────────────
def feature_card(x_mm, y_mm, icon, title, desc):
    w_mm, h_mm = 91, 50
    add_rect(s1, Mm(x_mm), Mm(y_mm), Mm(w_mm), Mm(h_mm),
             fill=WHITE, line=LINE, line_w=0.75)
    # アイコン背景（左上の角丸風）
    add_rect(s1, Mm(x_mm + 4), Mm(y_mm + 4), Mm(13), Mm(13), fill=LIGHT)
    add_text(s1, Mm(x_mm + 4), Mm(y_mm + 4), Mm(13), Mm(13), icon,
             size=20, color=BLUE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # タイトル + 説明
    add_text(s1, Mm(x_mm + 21), Mm(y_mm + 4), Mm(w_mm - 25), Mm(8), title,
             size=12.5, color=NAVY, bold=True, font=FONT_BOLD)
    add_text(s1, Mm(x_mm + 4), Mm(y_mm + 22), Mm(w_mm - 8), Mm(h_mm - 26), desc,
             size=9, color=INK2, font=FONT, line_spacing=1.45)

# 1段目
feature_card(12, 130, "▦", "16 分割 監視",
             "店舗選択で即座に 4×4 グリッド開始。\n"
             "Frigate API から 16 カメラの JPEG を並列取得し、\n"
             "サーバ側で合成して 5 秒間隔で配信。")
feature_card(107, 130, "●", "シングルライブ",
             "セルクリックで対象カメラだけを 1 秒間隔で配信。\n"
             "WebRTC 不要のため起動が速く、\n"
             "通信量も最小限。")

# 2段目
feature_card(12, 185, "◆", "VOD 録画再生",
             "Uniview / Frigate / i-PRO レコーダの\n"
             "録画にシーク可能。LiveKit Cloud Ingress で\n"
             "WebRTC 配信、画質と滑らかさを両立。")
feature_card(107, 185, "⚠", "BCP / SECURITY / INFRA",
             "J-Alert 自動クリップ、警備巡回スナップショット、\n"
             "死活監視・無人ストリーム検知まで一体化。\n"
             "PDF レポートを自動生成。")

# フッター帯
add_rect(s1, Mm(0), Mm(245), Mm(210), Mm(52), fill=LIGHT)
add_text(
    s1, Mm(12), Mm(250), Mm(186), Mm(8),
    "選ばれる 3 つの理由",
    size=13, color=NAVY, bold=True, font=FONT_BOLD,
)

def reason(x_mm, num, head, body):
    add_rect(s1, Mm(x_mm), Mm(263), Mm(11), Mm(11), fill=NAVY)
    add_text(s1, Mm(x_mm), Mm(263), Mm(11), Mm(11), num,
             size=14, color=WHITE, bold=True,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=FONT_BOLD)
    add_text(s1, Mm(x_mm + 14), Mm(262), Mm(50), Mm(6), head,
             size=10.5, color=NAVY, bold=True, font=FONT_BOLD)
    add_text(s1, Mm(x_mm + 14), Mm(269), Mm(50), Mm(20), body,
             size=8.5, color=INK2, font=FONT, line_spacing=1.4)

reason(12,  "1", "コスト 1 店舗 ¥58/月", "10,000 店舗運用時の月額\nクラウドコストを按分。")
reason(78,  "2", "ポート開放不要",       "店舗側は HTTPS 443\nアウトバウンドのみ。")
reason(144, "3", "30 分でデプロイ",      "Vercel + Supabase で\n運用人員 1 名を目標。")

# ページ番号
add_text(s1, Mm(180), Mm(289), Mm(20), Mm(6),
         "front / 1", size=7.5, color=MUTED, align=PP_ALIGN.RIGHT, font=FONT)


# ─────────────────────────────────────────────────────────────────────
# ▼ Slide 2 (裏面)
# ─────────────────────────────────────────────────────────────────────
s2 = add_blank_slide()

# 裏面ヘッダ（薄）
add_rect(s2, Mm(0), Mm(0), Mm(210), Mm(15), fill=NAVY)
add_text(s2, Mm(12), Mm(4), Mm(186), Mm(8),
         "RECORDER MONITOR  /  詳細仕様・料金・お問い合わせ",
         size=10, color=WHITE, bold=True, font=FONT_BOLD)

# セクション 1: 機能マトリクス
add_text(s2, Mm(12), Mm(20), Mm(186), Mm(8),
         "■ 機能一覧",
         size=12.5, color=NAVY, bold=True, font=FONT_BOLD)

def matrix_row(y_mm, name, monitor, mobile, note):
    add_rect(s2, Mm(12), Mm(y_mm), Mm(186), Mm(7),
             fill=None, line=LINE, line_w=0.3)
    add_text(s2, Mm(14), Mm(y_mm), Mm(60), Mm(7), name,
             size=8.5, color=INK, anchor=MSO_ANCHOR.MIDDLE, font=FONT)
    add_text(s2, Mm(75), Mm(y_mm), Mm(20), Mm(7), monitor,
             size=8.5, color=INK2, anchor=MSO_ANCHOR.MIDDLE,
             align=PP_ALIGN.CENTER, font=FONT)
    add_text(s2, Mm(96), Mm(y_mm), Mm(20), Mm(7), mobile,
             size=8.5, color=INK2, anchor=MSO_ANCHOR.MIDDLE,
             align=PP_ALIGN.CENTER, font=FONT)
    add_text(s2, Mm(117), Mm(y_mm), Mm(80), Mm(7), note,
             size=8, color=MUTED, anchor=MSO_ANCHOR.MIDDLE, font=FONT)

# ヘッダ行
add_rect(s2, Mm(12), Mm(30), Mm(186), Mm(7), fill=LIGHT, line=LINE, line_w=0.3)
add_text(s2, Mm(14),  Mm(30), Mm(60), Mm(7), "機能",
         size=9, color=NAVY, bold=True, anchor=MSO_ANCHOR.MIDDLE, font=FONT_BOLD)
add_text(s2, Mm(75),  Mm(30), Mm(20), Mm(7), "PC Web",
         size=9, color=NAVY, bold=True, anchor=MSO_ANCHOR.MIDDLE,
         align=PP_ALIGN.CENTER, font=FONT_BOLD)
add_text(s2, Mm(96),  Mm(30), Mm(20), Mm(7), "スマホ PWA",
         size=9, color=NAVY, bold=True, anchor=MSO_ANCHOR.MIDDLE,
         align=PP_ALIGN.CENTER, font=FONT_BOLD)
add_text(s2, Mm(117), Mm(30), Mm(80), Mm(7), "備考",
         size=9, color=NAVY, bold=True, anchor=MSO_ANCHOR.MIDDLE, font=FONT_BOLD)

rows = [
    ("16 分割監視（grid JPEG）",       "○", "○ 2×2",  "店舗選択で auto-start、5 秒更新"),
    ("単一カメラライブ",                "○", "○",      "1 秒間隔ポーリング、起動 1–2 秒"),
    ("VOD 録画再生",                   "○", "○",      "シークバー、最大 90 分セッション"),
    ("BCP（J-Alert クリップ録画）",    "○", "閲覧",   "ミサイル/地震/津波対応＋テスト発令"),
    ("SECURITY（警備巡回）",           "○", "閲覧",   "ベースライン差分 + AI（任意）"),
    ("INFRA 死活監視",                 "○", "○",      "edge 死活、無人ストリーム watchdog"),
    ("レポート PDF 自動生成",          "○", "閲覧",   "BCP/SECURITY/INFRA 各種"),
    ("マルチテナント + RBAC",          "○", "○",      "RLS による行レベル制御"),
    ("MFA 認証 + 監査ログ",            "○", "○",      "TOTP 必須、操作ログ 1 年保持"),
    ("プッシュ通知",                   "─", "○",      "Web Push（iOS 16.4+）"),
]
for i, (n, m, sp, nt) in enumerate(rows):
    matrix_row(37 + i * 7, n, m, sp, nt)

# セクション 2: アーキ概略
add_text(s2, Mm(12), Mm(115), Mm(95), Mm(8),
         "■ アーキテクチャ概略",
         size=12.5, color=NAVY, bold=True, font=FONT_BOLD)

add_rect(s2, Mm(12), Mm(125), Mm(95), Mm(55),
         fill=RGBColor(0x0F, 0x17, 0x2A), line=NAVY, line_w=0.5)
add_text(
    s2, Mm(14), Mm(126), Mm(91), Mm(54),
    "[現地店舗 × N]\n"
    "  カメラ / NVR / Frigate\n"
    "          │\n"
    "  HTTP JPEG / RTSP / MP4\n"
    "          ▼\n"
    "  Edge Agent (Node 22)\n"
    "  状態機械: Idle / Grid / Live / VOD / BCP\n"
    "          │ Outbound HTTPS 443\n"
    "          ▼\n"
    "[クラウド]\n"
    "  Vercel  (Next.js 16)\n"
    "  Supabase (Postgres + Auth + Storage)\n"
    "  LiveKit Cloud  ※ VOD のみ\n"
    "  Cloudflare (CDN + WAF)",
    size=8, color=RGBColor(0xE2, 0xE8, 0xF0), font="SF Mono", line_spacing=1.4,
)

# セクション 3: 料金（右側）
add_text(s2, Mm(112), Mm(115), Mm(86), Mm(8),
         "■ 料金（参考・税抜）",
         size=12.5, color=NAVY, bold=True, font=FONT_BOLD)

def price_row(y_mm, scale, cost_total, per_store):
    add_rect(s2, Mm(112), Mm(y_mm), Mm(86), Mm(8),
             fill=None, line=LINE, line_w=0.3)
    add_text(s2, Mm(114), Mm(y_mm), Mm(30), Mm(8), scale,
             size=9, color=INK, anchor=MSO_ANCHOR.MIDDLE, bold=True, font=FONT_BOLD)
    add_text(s2, Mm(144), Mm(y_mm), Mm(30), Mm(8), cost_total,
             size=9, color=INK2, anchor=MSO_ANCHOR.MIDDLE,
             align=PP_ALIGN.CENTER, font=FONT)
    add_text(s2, Mm(175), Mm(y_mm), Mm(22), Mm(8), per_store,
             size=9, color=BLUE, anchor=MSO_ANCHOR.MIDDLE,
             align=PP_ALIGN.RIGHT, bold=True, font=FONT_BOLD)

# 料金表ヘッダ
add_rect(s2, Mm(112), Mm(125), Mm(86), Mm(8), fill=LIGHT, line=LINE, line_w=0.3)
add_text(s2, Mm(114), Mm(125), Mm(30), Mm(8), "規模",
         size=8.5, color=NAVY, bold=True, anchor=MSO_ANCHOR.MIDDLE, font=FONT_BOLD)
add_text(s2, Mm(144), Mm(125), Mm(30), Mm(8), "月額合計",
         size=8.5, color=NAVY, bold=True, anchor=MSO_ANCHOR.MIDDLE,
         align=PP_ALIGN.CENTER, font=FONT_BOLD)
add_text(s2, Mm(175), Mm(125), Mm(22), Mm(8), "1 店舗",
         size=8.5, color=NAVY, bold=True, anchor=MSO_ANCHOR.MIDDLE,
         align=PP_ALIGN.RIGHT, font=FONT_BOLD)

price_row(133, "100 店舗",    "¥ 33,175",  "¥ 332")
price_row(141, "1,000 店舗",  "¥ 208,900", "¥ 209")
price_row(149, "10,000 店舗", "¥ 578,500", "¥ 58")

add_text(s2, Mm(112), Mm(160), Mm(86), Mm(20),
         "※ Vercel / Supabase / LiveKit / Cloudflare /\n"
         "観測 (Better Stack 等) の合算、為替 ¥155/$1。\n"
         "別途、エッジサーバ機材 ¥45,000/台前後。\n"
         "現地設置・LTE バックアップ等は別途お見積。",
         size=8, color=MUTED, font=FONT, line_spacing=1.45)

# セクション 4: 動作環境
add_text(s2, Mm(12), Mm(185), Mm(186), Mm(8),
         "■ 動作環境・対応機器",
         size=12.5, color=NAVY, bold=True, font=FONT_BOLD)

env_grid = [
    ("対応レコーダ", "i-PRO（WJ-NX410/510/NU101 ほか）\nUniview（NVR301/302/304 系）\nFrigate（OSS + go2rtc 構成）"),
    ("クライアント", "デスクトップ: Chrome / Safari 最新版\nモバイル: iOS 16.4+ / Android 12+\nPWA インストール対応（ホーム画面追加）"),
    ("クラウド",     "Vercel (Pro/Enterprise)\nSupabase (Team/Enterprise)\nLiveKit Cloud (Scale)"),
    ("ネットワーク", "店舗側: アウトバウンド HTTPS 443 のみ\nUPS / LTE バックアップ任意（推奨）\nプロキシ環境 対応"),
]
for i, (h, b) in enumerate(env_grid):
    x = 12 + (i % 2) * 95
    y = 195 + (i // 2) * 25
    add_rect(s2, Mm(x), Mm(y), Mm(91), Mm(22), fill=WHITE, line=LINE, line_w=0.5)
    add_text(s2, Mm(x + 3), Mm(y + 2), Mm(85), Mm(5), h,
             size=9, color=NAVY, bold=True, font=FONT_BOLD)
    add_text(s2, Mm(x + 3), Mm(y + 8), Mm(85), Mm(13), b,
             size=8, color=INK2, font=FONT, line_spacing=1.4)

# セクション 5: お問い合わせ（フッター）
add_rect(s2, Mm(0), Mm(255), Mm(210), Mm(42), fill=NAVY)
add_text(s2, Mm(12), Mm(260), Mm(186), Mm(8),
         "お問い合わせ・無料 PoC（パイロット 3 店舗）",
         size=14, color=WHITE, bold=True, font=FONT_BOLD)
add_text(
    s2, Mm(12), Mm(270), Mm(186), Mm(20),
    "3 店舗規模のパイロット（1 か月）から開始可能です。\n"
    "実機 / クラウド / UX のすべてを現場でご確認いただいた上で、本展開へ。\n\n"
    "TEL  03-0000-0000      MAIL  sales@example.com      WEB  https://example.com/recorder",
    size=9.5, color=RGBColor(0xDB, 0xEA, 0xFE), font=FONT, line_spacing=1.6,
)

# ページ番号
add_text(s2, Mm(180), Mm(289), Mm(20), Mm(6),
         "back / 2", size=7.5, color=RGBColor(0xCB, 0xD5, 0xE1),
         align=PP_ALIGN.RIGHT, font=FONT)

# 保存
prs.save(OUT)
print(f"saved: {OUT}")
